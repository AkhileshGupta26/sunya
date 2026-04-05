from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent.parent
TEMPLATE = Path(r"C:\Users\MSI\Downloads\69cbc185a9372_HACDAYS_PPT_TEMPLATE.pptx")
OUTPUT = ROOT / "Sunya_Hackathon_Showcase_Deck.pptx"

BG = RGBColor(7, 18, 33)
CARD = RGBColor(15, 34, 54)
CARD_ALT = RGBColor(22, 49, 79)
ACCENT = RGBColor(45, 212, 191)
ACCENT_2 = RGBColor(251, 191, 36)
TEXT = RGBColor(241, 245, 249)
MUTED = RGBColor(191, 219, 254)
RED = RGBColor(248, 113, 113)
GREEN = RGBColor(74, 222, 128)


def clear_body(slide):
    for shape in list(slide.shapes):
        if shape.shape_type == 14:
            continue
        if hasattr(shape, "text") and shape.text:
            shape.text = ""


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def style_text_frame(
    text_frame,
    text,
    font_size=18,
    bold=False,
    color=TEXT,
    font_name="Aptos",
    align=PP_ALIGN.LEFT,
):
    text_frame.clear()
    lines = text.split("\n")
    for idx, line in enumerate(lines):
        p = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        p.text = line
        p.alignment = align
        for run in p.runs:
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.name = font_name
            run.font.color.rgb = color


def add_textbox(slide, left, top, width, height, text, **kwargs):
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    style_text_frame(shape.text_frame, text, **kwargs)
    shape.text_frame.word_wrap = True
    shape.text_frame.vertical_anchor = MSO_ANCHOR.TOP
    return shape


def add_round_box(slide, left, top, width, height, fill_color, line_color=None, radius=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color or fill_color
    shape.line.width = Pt(1.5)
    return shape


def add_labeled_card(slide, left, top, width, height, title, body, fill_color=CARD):
    box = add_round_box(slide, left, top, width, height, fill_color, CARD_ALT)
    tf = box.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    p1.text = title
    p1.alignment = PP_ALIGN.LEFT
    for run in p1.runs:
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.color.rgb = TEXT
        run.font.name = "Aptos"
    p2 = tf.add_paragraph()
    p2.text = body
    for run in p2.runs:
        run.font.size = Pt(12)
        run.font.color.rgb = MUTED
        run.font.name = "Aptos"
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)
    tf.margin_top = Pt(8)
    return box


def add_connector(slide, x1, y1, x2, y2, color=ACCENT):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = color
    line.line.width = Pt(2.2)
    return line


def duplicate_template(prs):
    layout = prs.slide_layouts[6]
    slides = []
    for idx in range(7):
        source = Presentation(str(TEMPLATE)).slides[idx]
        slide = prs.slides.add_slide(layout)
        for shape in source.shapes:
            el = deepcopy(shape.element)
            slide.shapes._spTree.insert_element_before(el, "p:extLst")
        slides.append(slide)
    return slides


def build_title_slide(slide):
    clear_body(slide)
    set_slide_bg(slide, BG)
    add_round_box(slide, 0.35, 0.95, 4.2, 2.8, CARD_ALT)
    add_textbox(slide, 0.65, 1.25, 3.6, 0.55, "SUNYA", font_size=28, bold=True, color=ACCENT)
    add_textbox(
        slide,
        0.65,
        1.8,
        3.4,
        0.95,
        "AI-powered mindful routine coach for digital detox, meditation, and habit streaks.",
        font_size=19,
        bold=True,
    )
    add_textbox(
        slide,
        0.65,
        2.75,
        3.4,
        0.55,
        "From scattered attention to calm, measurable daily rituals.",
        font_size=13,
        color=MUTED,
    )
    add_labeled_card(slide, 4.9, 1.15, 4.45, 0.9, "Team", "Sunya Labs\nLead builder: Akhil Smile")
    add_labeled_card(slide, 4.9, 2.2, 4.45, 0.9, "Hackathon Tracks", "Healthcare | AI/ML | Open Innovation")
    add_labeled_card(slide, 4.9, 3.25, 4.45, 0.9, "MVP Status", "Working Expo frontend + FastAPI backend + MongoDB + Gemini-powered Sunya Yogi")
    add_textbox(slide, 0.6, 4.32, 7.2, 0.35, "Startup showcase deck built from the current MVP repository.", font_size=11, color=MUTED)


def build_problem_solution(slide):
    clear_body(slide)
    set_slide_bg(slide, BG)
    add_textbox(slide, 0.45, 0.75, 2.0, 0.3, "THE PROBLEM", font_size=18, bold=True, color=RED)
    add_textbox(slide, 5.1, 0.75, 2.6, 0.3, "THE SUNYA SOLUTION", font_size=18, bold=True, color=GREEN)

    problems = [
        ("Attention fatigue", "Young users want calm, but apps often increase screen time instead of reducing it."),
        ("No accountability loop", "Meditation products rarely build habit streaks, social motivation, and measurable progress together."),
        ("Generic wellness advice", "Most tools do not personalize routines using emotion, context, and daily behavior."),
    ]
    y = 1.2
    for title, body in problems:
        add_labeled_card(slide, 0.45, y, 4.0, 0.95, title, body, fill_color=CARD)
        y += 1.07

    solutions = [
        "Guided meditation + yoga + digital detox in one daily flow.",
        "Sunya Yogi AI recommends the right session based on how the user feels.",
        "Points, streaks, leaderboards, and contests keep consistency rewarding.",
        "Optional BPM-assisted mindfulness check adds a future health-tech edge.",
    ]
    box = add_round_box(slide, 5.0, 1.18, 4.35, 2.7, CARD, CARD_ALT)
    tf = box.text_frame
    tf.clear()
    for idx, line in enumerate(solutions):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"- {line}"
        for run in p.runs:
            run.font.size = Pt(15)
            run.font.color.rgb = TEXT
            run.font.name = "Aptos"
    tf.word_wrap = True
    add_labeled_card(
        slide,
        5.0,
        4.02,
        4.35,
        0.7,
        "MVP Outcome",
        "Users complete a calming session, earn points, enter detox mode, and track progress in one app.",
        fill_color=CARD_ALT,
    )


def build_tech_slide(slide):
    clear_body(slide)
    set_slide_bg(slide, BG)
    columns = [
        (0.42, "Frontend", "Expo Router\nReact Native\nTypeScript\nReanimated UI"),
        (2.85, "Backend", "FastAPI\nPython\nJWT auth\nMotor + MongoDB"),
        (5.28, "AI + APIs", "Gemini API\nCustom Yogi prompt\nCamera/BPM flow\nNotifications"),
        (7.4, "Product Layer", "Streaks\nContests\nLeaderboards\nRoutine likes"),
    ]
    for left, title, body in columns:
        add_labeled_card(slide, left, 1.25, 2.05, 2.55, title, body, fill_color=CARD if left < 5.3 else CARD_ALT)
    add_round_box(slide, 0.55, 4.15, 8.9, 0.75, CARD_ALT, ACCENT)
    add_textbox(
        slide,
        0.8,
        4.32,
        8.3,
        0.35,
        "Current MVP already supports authentication, daily sessions, progress history, detox scoring, contest ranking, and AI-guided Yogi recommendations.",
        font_size=13,
        bold=True,
    )


def build_workflow_slide(slide):
    clear_body(slide)
    set_slide_bg(slide, BG)
    steps = [
        ("1. Onboard", "Create account\nSet wake habit\nChoose intent"),
        ("2. Daily prompt", "Open home\nSee focus card\nStart routine"),
        ("3. Guided action", "Meditation / Yoga\nOptional BPM\nMood-aware Yogi"),
        ("4. Detox lock-in", "Start 30-120 min detox\nStay present\nEarn points"),
        ("5. Progress loop", "Track streaks\nJoin contests\nReturn tomorrow"),
    ]
    lefts = [0.3, 2.15, 4.0, 5.85, 7.7]
    for idx, (title, body) in enumerate(steps):
        add_labeled_card(slide, lefts[idx], 1.8, 1.45, 1.65, title, body, fill_color=CARD if idx % 2 == 0 else CARD_ALT)
        if idx < len(steps) - 1:
            add_connector(slide, lefts[idx] + 1.45, 2.62, lefts[idx + 1], 2.62)
    add_textbox(
        slide,
        0.8,
        4.12,
        8.2,
        0.45,
        "This workflow converts mindfulness from a one-time session into a sticky habit loop with emotional personalization and gamified retention.",
        font_size=13,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )


def build_architecture_slide(slide):
    clear_body(slide)
    set_slide_bg(slide, BG)
    add_labeled_card(slide, 0.35, 1.4, 2.0, 2.0, "Client App", "Expo app\nAuth\nHome dashboard\nMeditation, Detox, Progress", fill_color=CARD_ALT)
    add_labeled_card(slide, 2.8, 1.05, 2.0, 0.95, "API Layer", "FastAPI routers\nJWT auth\nCORS", fill_color=CARD)
    add_labeled_card(slide, 2.8, 2.25, 2.0, 1.15, "Core Services", "Sessions\nUsers\nContests\nRoutines\nYogi", fill_color=CARD)
    add_labeled_card(slide, 5.25, 1.05, 2.0, 0.95, "Data Store", "MongoDB Atlas\nUsers\nSessions\nLikes", fill_color=CARD)
    add_labeled_card(slide, 5.25, 2.25, 2.0, 1.15, "AI Layer", "Gemini-powered Sunya Yogi\nTrack recommendation\nReflection output", fill_color=CARD_ALT)
    add_labeled_card(slide, 7.7, 1.4, 1.6, 2.0, "External Signals", "Camera/BPM\nPush alerts\nWake-time reminders", fill_color=CARD)

    add_connector(slide, 2.35, 2.0, 2.8, 2.0)
    add_connector(slide, 4.8, 1.52, 5.25, 1.52)
    add_connector(slide, 4.8, 2.8, 5.25, 2.8)
    add_connector(slide, 7.25, 2.0, 7.7, 2.0)
    add_textbox(slide, 0.8, 4.15, 8.1, 0.4, "Modular backend means each growth feature can scale independently without rewriting the product core.", font_size=13, color=MUTED, align=PP_ALIGN.CENTER)


def build_usp_slide(slide):
    clear_body(slide)
    set_slide_bg(slide, BG)
    cards = [
        (0.55, 1.45, "Not another meditation timer", "Sunya blends meditation, yoga, detox, and social motivation in one experience."),
        (5.05, 1.45, "Personalization that feels human", "Sunya Yogi converts a plain text check-in into a recommended session and emotional reflection."),
        (0.55, 3.0, "Retention by design", "Streaks, points, contests, and leaderboards make calm behavior repeatable."),
        (5.05, 3.0, "Low-cost, scalable wellness", "Software-first model works for students, wellness communities, and employers without heavy hardware."),
    ]
    for left, top, title, body in cards:
        add_labeled_card(slide, left, top, 4.0, 1.1, title, body, fill_color=CARD if left < 1 else CARD_ALT)


def build_feasibility_slide(slide):
    clear_body(slide)
    set_slide_bg(slide, BG)
    sections = [
        ("Why feasible now", "MVP already exists in code: working frontend, backend APIs, auth, progress tracking, contests, and AI coaching."),
        ("User wedge", "Students and young professionals facing burnout, doomscrolling, and inconsistent self-care."),
        ("Revenue path", "Freemium app, premium guided programs, campus wellness partnerships, and B2B employee wellbeing packs."),
        ("Next 30 days", "Polish onboarding, re-enable accurate BPM, launch beta cohort, measure D7 retention and daily streak completion."),
    ]
    y = 1.25
    for title, body in sections:
        add_labeled_card(slide, 0.7, y, 8.6, 0.72, title, body, fill_color=CARD if y < 2.8 else CARD_ALT)
        y += 0.86
    add_textbox(slide, 2.1, 4.76, 5.8, 0.25, "Sunya is credible as a hackathon startup because the product loop is already implemented, not just imagined.", font_size=11, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    slides = duplicate_template(prs)

    build_title_slide(slides[0])
    build_problem_solution(slides[1])
    build_tech_slide(slides[2])
    build_workflow_slide(slides[3])
    build_architecture_slide(slides[4])
    build_usp_slide(slides[5])
    build_feasibility_slide(slides[6])

    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    main()
